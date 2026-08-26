#!/usr/bin/env python3
# =============================================================================
#  Into Frame -- slide deck generator
#
#  Builds Into-Frame-slides.pptx (20 slides, 16:9) from the content of the
#  report in sections/, using the figures in figures/.
#
#  Run:  python3 make_slides.py            (from report/)
#  Deps: python-pptx, Pillow, PyMuPDF (the last only for the Paris height-field
#        figure, which is a PDF: it is rasterised and split into its two panels
#        so the slide can carry its own labels on the dark background).
#
#  Everything below the helpers is content. To change wording, edit build();
#  to change the look, edit the palette and the helpers.
#
#  Note on glyphs: Helvetica Neue has no arrow or subscript characters, so the
#  deck spells those out (`sub()` for D_max, words instead of arrows). Do not
#  reintroduce them.
# =============================================================================
import os
import tempfile

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")
OUT = os.path.join(HERE, "Into-Frame-slides.pptx")
TMP = tempfile.mkdtemp(prefix="intoframe-slides-")

# ---- canvas -----------------------------------------------------------------
W, H = 13.333, 7.5
ML, MR = 0.72, 0.72
CONTENT_W = W - ML - MR
TOP_RULE = 1.46
BODY_TOP = 1.72
BODY_BOT = 6.86

# ---- palette ----------------------------------------------------------------
BG      = RGBColor(0x0E, 0x11, 0x16)
PANEL   = RGBColor(0x17, 0x1C, 0x24)
PANEL2  = RGBColor(0x1E, 0x24, 0x2E)
HAIR    = RGBColor(0x2C, 0x34, 0x40)
FAINT   = RGBColor(0x1E, 0x24, 0x2E)
INK     = RGBColor(0xF2, 0xF4, 0xF7)
BODY    = RGBColor(0xC2, 0xCA, 0xD5)
MUTED   = RGBColor(0x7B, 0x87, 0x97)
DIM     = RGBColor(0x5A, 0x64, 0x72)
TEAL    = RGBColor(0x6F, 0xD2, 0xC0)
AMBER   = RGBColor(0xE9, 0xA9, 0x4B)
CLAY    = RGBColor(0xDE, 0x73, 0x58)
SLATE   = RGBColor(0x93, 0xA4, 0xB8)
BLUE    = RGBColor(0x6E, 0xA8, 0xE0)
GREEN   = RGBColor(0x8F, 0xC9, 0x6B)

FONT = "Helvetica Neue"

TRACK_COLORS = {
    "Terrain": TEAL, "Objects": AMBER, "Ground cover": GREEN,
    "Water": BLUE, "Sky": RGBColor(0xD8, 0xC0, 0x6A),
    "Motion": RGBColor(0xB7, 0x9B, 0xE0),
}


# =============================================================================
#  drawing helpers
# =============================================================================
def _alpha(clr_elm, pct):
    a = clr_elm.makeelement(qn("a:alpha"), {})
    a.set("val", str(int(pct * 1000)))
    clr_elm.append(a)


def rect(slide, x, y, w, h, color, alpha=None, line=None, line_w=0.75,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if alpha is not None:
        _alpha(s.fill._xPr.find(qn("a:solidFill")).find(qn("a:srgbClr")), alpha)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = radius
    s.text_frame.word_wrap = True
    return s


def gradient_scrim(slide, x, y, w, h, a_top, a_bottom, color=BG):
    """A vertical fade, used to seat title text over a photograph."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.shadow.inherit = False
    s.line.fill.background()
    s.fill.gradient()
    for stop in s.fill.gradient_stops:
        stop.color.rgb = color
    grad = s.fill._xPr.find(qn("a:gradFill"))
    for gs, a in zip(grad.find(qn("a:gsLst")).findall(qn("a:gs")),
                     (a_top, a_bottom)):
        _alpha(gs.find(qn("a:srgbClr")), a)
    # 5400000 = 90 degrees clockwise from east, i.e. top to bottom
    grad.find(qn("a:lin")).set("ang", "5400000")
    return s


def line(slide, x1, y1, x2, y2, color=HAIR, width=0.75):
    from pptx.enum.shapes import MSO_CONNECTOR
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def style(run, size, color=BODY, bold=False, italic=False, spacing=None,
          baseline=None):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT
    if spacing is not None:
        run.font._rPr.set("spc", str(int(spacing)))
    if baseline is not None:
        run.font._rPr.set("baseline", str(int(baseline)))
    return run


def para(tf, first=False):
    return tf.paragraphs[0] if first else tf.add_paragraph()


def rich(p, text, size, color=BODY, bold_color=INK, **kw):
    """Write text into p, honouring **bold**, *italic* and D{max} subscripts.
    Emphasis markers do not nest in each other, but a subscript may appear
    inside one."""
    import re

    def emit(txt, col, bold=False, italic=False):
        for bit in re.split(r"(\{[a-z]+\})", txt):
            if not bit:
                continue
            if bit.startswith("{"):
                style(p.add_run(), size * 0.78, col, bold=bold, italic=italic,
                      baseline=-25000, **kw).text = bit[1:-1]
            else:
                style(p.add_run(), size, col, bold=bold, italic=italic,
                      **kw).text = bit

    for chunk in re.split(r"(\*\*[^*]+\*\*|\*[^*]+\*)", text):
        if not chunk:
            continue
        if chunk.startswith("**"):
            emit(chunk[2:-2], bold_color, bold=True)
        elif chunk.startswith("*"):
            emit(chunk[1:-1], color, italic=True)
        else:
            emit(chunk, color)
    return p


def hanging(p, indent=0.19):
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(Emu(int(indent * 914400))))
    pPr.set("indent", str(-Emu(int(indent * 914400))))


def bullets(slide, x, y, w, items, size=12.5, gap=10, dash_color=TEAL,
            line_spacing=1.22, h=None):
    tf = textbox(slide, x, y, w, h or (BODY_BOT - y))
    for i, item in enumerate(items):
        p = para(tf, first=(i == 0))
        p.space_after = Pt(gap)
        p.line_spacing = line_spacing
        hanging(p)
        style(p.add_run(), size, dash_color, bold=True).text = "— "
        rich(p, item, size)
    return tf


def add_pic(slide, path, x, y, w, h, mode="cover"):
    iw, ih = Image.open(path).size
    a, b = iw / ih, w / h
    if mode == "contain":
        if a > b:
            h2 = w / a
            y += (h - h2) / 2
            h = h2
        else:
            w2 = h * a
            x += (w - w2) / 2
            w = w2
        return slide.shapes.add_picture(path, Inches(x), Inches(y),
                                        Inches(w), Inches(h))
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    if a > b:
        f = (1 - b / a) / 2
        pic.crop_left = pic.crop_right = f
    else:
        f = (1 - a / b) / 2
        pic.crop_top = pic.crop_bottom = f
    return pic


def zoom_to_content(path, pad=0.10, thresh=24):
    """Crop an image to the bounding box of its non-black content, so a small
    mask on a large empty grid does not read as an empty box."""
    im = Image.open(path).convert("RGB")
    bbox = im.convert("L").point(lambda v: 255 if v > thresh else 0).getbbox()
    if not bbox:
        return path
    x0, y0, x1, y1 = bbox
    px, py = int((x1 - x0) * pad), int((y1 - y0) * pad)
    side = max(x1 - x0 + 2 * px, y1 - y0 + 2 * py)
    cx, cy = (x0 + x1) // 2, (y0 + y1) // 2
    box = (cx - side // 2, cy - side // 2, cx + side // 2, cy + side // 2)
    out = os.path.join(TMP, "zoom-" + os.path.basename(path))
    im.crop(box).save(out)
    return out


def caption(slide, x, y, w, text, size=8.5, color=MUTED, align=PP_ALIGN.LEFT,
            h=0.6):
    tf = textbox(slide, x, y, w, h)
    p = para(tf, first=True)
    p.alignment = align
    p.line_spacing = 1.18
    rich(p, text, size, color)
    return tf


def chip(slide, x, y, w, h, label, color=TEAL, size=9.5, fill_alpha=13):
    s = rect(slide, x, y, w, h, color, alpha=fill_alpha,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.28)
    tf = s.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.06)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    style(p.add_run(), size, color, bold=True).text = label
    return s


def eyebrow(slide, x, y, w, text, color=TEAL):
    tf = textbox(slide, x, y, w, 0.26)
    style(tf.paragraphs[0].add_run(), 9.5, color, bold=True,
          spacing=140).text = text.upper()
    return tf


def card(slide, x, y, w, h, accent, head, body, head_size=13, body_size=10.5,
         bar="top"):
    rect(slide, x, y, w, h, PANEL, radius=0.07,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if bar == "top":
        rect(slide, x, y, w, 0.035, accent)
    else:
        rect(slide, x, y, 0.045, h, accent)
    pad = 0.32 if bar == "top" else 0.36
    tf = textbox(slide, x + pad, y + 0.16, w - pad - 0.28, h - 0.32,
                 anchor=MSO_ANCHOR.MIDDLE)
    style(tf.paragraphs[0].add_run(), head_size, accent, bold=True).text = head
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    p.line_spacing = 1.22
    rich(p, body, body_size, BODY, bold_color=INK)
    return tf


# =============================================================================
#  deck chrome
# =============================================================================
class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(W)
        self.prs.slide_height = Inches(H)
        self.blank = self.prs.slide_layouts[6]
        self.n = 0

    def slide(self, kick=None, title=None, footer=None, rule=True):
        s = self.prs.slides.add_slide(self.blank)
        self.n += 1
        rect(s, 0, 0, W, H, BG)
        if kick:
            eyebrow(s, ML, 0.50, CONTENT_W, kick)
        if title:
            tf = textbox(s, ML, 0.76, CONTENT_W, 0.62)
            p = tf.paragraphs[0]
            p.line_spacing = 1.0
            rich(p, title, 27, INK, spacing=-8)
        if rule:
            line(s, ML, TOP_RULE, W - MR, TOP_RULE, HAIR, 0.75)
        if self.n > 1:
            tf = textbox(s, W - MR - 0.7, 6.98, 0.7, 0.24)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            style(p.add_run(), 9, MUTED).text = str(self.n)
            if footer:
                tf = textbox(s, ML, 6.98, 6.0, 0.24)
                style(tf.paragraphs[0].add_run(), 9, DIM, spacing=60).text = footer
        return s


# =============================================================================
#  derived figures
# =============================================================================
def paris_panels():
    src = os.path.join(FIG, "paris-terrain.pdf")
    try:
        import fitz
    except ImportError:
        print("  ! PyMuPDF missing -- skipping the Paris height-field panels")
        return None, None
    pix = fitz.open(src)[0].get_pixmap(dpi=260)
    raw = os.path.join(TMP, "paris-terrain.png")
    pix.save(raw)
    im = Image.open(raw).convert("RGB")
    g = im.convert("L")
    wd, ht = g.size
    def spans(flags, min_len, max_gap):
        """Contiguous runs of True, bridging gaps shorter than max_gap. The
        panels contain bright pixels, so a strict run split them in three."""
        out, start = [], None
        for i, f in enumerate(list(flags) + [False]):
            if f and start is None:
                start = i
            elif not f and start is not None:
                if out and start - out[-1][1] <= max_gap:
                    out[-1] = (out[-1][0], i)
                else:
                    out.append((start, i))
                start = None
        return [r for r in out if r[1] - r[0] > min_len]

    ys = range(0, ht, 4)
    dark_col = [sum(g.getpixel((x, y)) < 110 for y in ys) > 0.5 * len(ys)
                for x in range(wd)]
    runs = spans(dark_col, wd * 0.15, wd * 0.04)
    if len(runs) != 2:
        print(f"  ! expected 2 panels in paris-terrain.pdf, found {len(runs)}")
        return None, None
    x0, x1 = runs[0]
    xs = range(x0, x1, 4)
    dark_row = [sum(g.getpixel((x, y)) < 110 for x in xs) > 0.6 * len(xs)
                for y in range(ht)]
    rows = spans(dark_row, ht * 0.2, ht * 0.04)
    y0, y1 = rows[0]
    out = []
    for i, (a, b) in enumerate(runs):
        p = os.path.join(TMP, f"paris-panel-{i}.png")
        im.crop((a, y0, b, y1)).save(p)
        out.append(p)
    return out


# =============================================================================
#  the deck
#
#  Copy is deliberately terse: a slide carries the claim and the number, the
#  speaker carries the sentence. If a bullet runs past two lines, cut it.
# =============================================================================
BULLET = 14      # body bullets
CARDT  = 12      # text inside a card
NOTE   = 12      # a full-width note under a rule
CAP    = 9       # figure captions


def build():
    d = Deck()
    F = lambda *p: os.path.join(FIG, *p)

    # -- 1 -- title ----------------------------------------------------------
    s = d.slide(rule=False)
    add_pic(s, F("rainier", "panorama.jpg"), 0, 0, W, H)
    gradient_scrim(s, 0, 0, W, 4.3, 90, 58)
    gradient_scrim(s, 0, 4.3, W, 3.2, 58, 28)
    tf = textbox(s, ML + 0.1, 1.30, 11.2, 2.2)
    p = tf.paragraphs[0]
    p.line_spacing = 1.02
    style(p.add_run(), 52, INK, bold=True, spacing=-30).text = "Jumping Into the Frame"
    p = tf.add_paragraph()
    p.space_before = Pt(14)
    p.line_spacing = 1.28
    style(p.add_run(), 17, RGBColor(0xCE, 0xD6, 0xE0)).text = (
        "Reconstructing interactive, physically-grounded 3D environments "
        "from a single image")
    line(s, ML + 0.1, 2.98, ML + 1.5, 2.98, TEAL, 1.6)
    tf = textbox(s, ML + 0.1, 3.22, 9.0, 0.4)
    p = tf.paragraphs[0]
    style(p.add_run(), 12.5, INK, bold=True).text = "Joshua Ford"
    style(p.add_run(), 12.5, MUTED).text = "        Advisor: "
    style(p.add_run(), 12.5, BODY).text = "Marie-Paule Cani"

    # -- 2 -- motivation -----------------------------------------------------
    s = d.slide("Motivation", "A photograph is not a place", "Introduction")
    bullets(s, ML, BODY_TOP, 6.4, [
        "One view, infinitely many scenes. The problem is underconstrained.",
        "People read a photograph as a place anyway: ground past the frame, the "
        "far side of a rock, individual flowers.",
        "Photographs are abundant. 3D authoring is not.",
        "The goal is somewhere you can **jump in**, not a picture you can orbit.",
    ], size=BULLET, gap=22)
    add_pic(s, F("rainier", "input.jpg"), 7.62, BODY_TOP, 4.99, 3.35)
    caption(s, 7.62, BODY_TOP + 3.45, 4.99,
            "One photograph — roughly a sixth of the sphere.", size=CAP)
    goals = [("Interactive", "objects animate and collide"),
             ("Decomposed", "entities, not painted texture"),
             ("Plausible", "beyond what was observed"),
             ("Performant", "stereo at headset frame rates")]
    for i, (g_, n_) in enumerate(goals):
        x = ML + i * 3.19
        line(s, x, 5.62, x + 2.95, 5.62, HAIR)
        tf = textbox(s, x, 5.78, 2.95, 0.9)
        style(tf.paragraphs[0].add_run(), 13, AMBER, bold=True).text = g_
        p = tf.add_paragraph()
        p.space_before = Pt(5)
        p.line_spacing = 1.2
        style(p.add_run(), 10.5, MUTED).text = n_

    # -- 3 -- the challenge ---------------------------------------------------
    s = d.slide("The challenge", "Four problems that compound", "Introduction")
    quad = [
        ("Extreme ill-posedness", TEAL,
         "**Five sixths of the sphere** is never observed. Nor is anything behind "
         "an object."),
        ("Compounding error", AMBER,
         "Depth, then geometry, then materials, then objects. Each inherits every "
         "earlier error."),
        ("Heterogeneous content", GREEN,
         "A mountainside, a tree, a meadow, the sea. Four representations, one "
         "coordinate frame, one light."),
        ("Interactivity budget", BLUE,
         "Stereo at headset frame rates — binding hardest on exactly the parts "
         "that move."),
    ]
    for i, (head, col, body) in enumerate(quad):
        card(s, ML + (i % 2) * 6.0, 2.10 + (i // 2) * 2.34, 5.62, 2.06,
             col, head, body, head_size=17, body_size=13.5, bar="left")

    # -- 4 -- related work ----------------------------------------------------
    s = d.slide("Related work", "Prior work stops at viewable", "Related work")
    groups = [
        ("Single-image panoramas",
         "CubeDiff · DreamCube · PanoDreamer · LayerPano3D"),
        ("Single-image 3D scenes",
         "SphericalDreamer · SonoWorld · 360Anything · WorldGen"),
        ("Terrain, texture, arrangement",
         "Pixels2Peaks · Earthbender · WorldBrush · Landlab · pattern texturing · "
         "point processes"),
        ("Perception and reconstruction",
         "Depth Anything 3 · SAM 2 · Grounding DINO · CLIP · SAM 3D · ObjectClear "
         "· LuxDiT · LTX-2"),
    ]
    y = BODY_TOP + 0.10
    for head, body in groups:
        line(s, ML, y - 0.16, 6.9, y - 0.16, FAINT)
        tf = textbox(s, ML, y, 6.5, 1.2)
        style(tf.paragraphs[0].add_run(), 14, AMBER, bold=True).text = head
        p = tf.add_paragraph()
        p.space_before = Pt(7)
        p.line_spacing = 1.24
        rich(p, body, 12, BODY)
        y += 1.22
    rect(s, 7.55, BODY_TOP - 0.06, 5.06, 4.60, PANEL, radius=0.05,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = textbox(s, 7.94, BODY_TOP + 0.14, 4.3, 4.2, anchor=MSO_ANCHOR.MIDDLE)
    style(tf.paragraphs[0].add_run(), 10, TEAL, bold=True,
          spacing=140).text = "WHERE WE DIFFER"
    p = tf.add_paragraph()
    p.space_before = Pt(14)
    p.line_spacing = 1.3
    rich(p, "They optimise fidelity near the original camera.\nWe optimise "
         "**physical decomposition** — typed entities, not one geometry proxy:",
         13.5, BODY)
    for item in ["a collidable terrain mesh",
                 "individually placed object meshes",
                 "an instanced ground-cover population",
                 "an animated water surface"]:
        p = tf.add_paragraph()
        p.space_before = Pt(9)
        hanging(p, 0.16)
        style(p.add_run(), 12.5, TEAL, bold=True).text = "· "
        style(p.add_run(), 12.5, BODY).text = item
    p = tf.add_paragraph()
    p.space_before = Pt(15)
    p.line_spacing = 1.28
    rich(p, "The panorama is an intermediate representation, not the deliverable.",
         12, MUTED)

    # -- 5 -- design principle ------------------------------------------------
    s = d.slide("Design principle", "Reconstruct each part by its physical role",
                "Overview")
    tf = textbox(s, ML, BODY_TOP, 11.9, 0.6)
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    rich(p, "**A scene is not one thing.** Ground, objects, ground cover, water "
         "and sky differ in how they are built, textured, and how they behave.",
         16, BODY)
    roles = [
        ("Terrain", "one height field, meshed at variable density, collidable"),
        ("Objects", "clustered by appearance, reconstructed once, instanced"),
        ("Ground cover", "scattered by density over a semantic and colour mask"),
        ("Water", "a separate levelled surface with an animated shader"),
        ("Sky", "skybox and image-based lighting source"),
        ("Motion", "a rigged sway, or a rigid body handed to physics"),
    ]
    for i, (name, body) in enumerate(roles):
        card(s, ML + (i % 3) * 4.03, 2.74 + (i // 3) * 1.72, 3.79, 1.52,
             TRACK_COLORS[name], name, body, head_size=15, body_size=11.5)
    tf = textbox(s, ML, 6.30, 11.9, 0.4)
    rich(tf.paragraphs[0], "One representation for all of them means accepting "
         "that representation's limits.", 12, MUTED)

    # -- 6 -- pipeline --------------------------------------------------------
    s = d.slide("42 stages, 38 active in the default configuration",
                "The pipeline", "Overview")
    pipeline_diagram(s)

    # -- 7 -- three panoramas -------------------------------------------------
    s = d.slide("The panoramic substrate",
                "Three panoramas, because they answer different questions",
                "The panoramic data")
    add_pic(s, F("rainier", "panorama-band.jpg"), ML, BODY_TOP, 7.4, 1.34)
    add_pic(s, F("rainier", "panorama-object-removed-band.jpg"),
            ML, BODY_TOP + 1.46, 7.4, 1.34)
    caption(s, ML, BODY_TOP + 2.88, 7.4,
            "Original (top) and object-removed (bottom): the meadow and the trees "
            "are gone with their shadows, exposing the ground behind them.",
            size=CAP)
    items = [
        ("Original", TEAL,
         "The only variant where the real objects remain. Every semantic stage "
         "reads it."),
        ("Object-removed", AMBER,
         "Depth through an occluder puts the flower canopy, not the ground, at "
         "the bottom of the height field."),
        ("Sky-inpainted", BLUE, "Skybox and image-based lighting."),
    ]
    y = BODY_TOP - 0.04
    for name, col, body in items:
        line(s, 8.44, y - 0.16, W - MR, y - 0.16, FAINT)
        tf = textbox(s, 8.44, y, 4.17, 1.4)
        style(tf.paragraphs[0].add_run(), 14.5, col, bold=True).text = name
        p = tf.add_paragraph()
        p.space_before = Pt(7)
        p.line_spacing = 1.26
        style(p.add_run(), 12, BODY).text = body
        y += 1.44
    line(s, ML, 5.32, W - MR, 5.32, HAIR)
    facts = [("~83%", "of the sphere is generated"),
             ("3", "panoramas kept in parallel"),
             ("2×", "region typing is run twice")]
    for i, (v, lab) in enumerate(facts):
        x = ML + i * 4.25
        tf = textbox(s, x, 5.56, 3.9, 1.1)
        p = tf.paragraphs[0]
        p.line_spacing = 1.0
        style(p.add_run(), 30, AMBER, bold=True, spacing=-20).text = v
        p = tf.add_paragraph()
        p.space_before = Pt(7)
        p.line_spacing = 1.2
        style(p.add_run(), 11.5, MUTED).text = lab

    # -- 8 -- calibration -----------------------------------------------------
    s = d.slide("Metric calibration",
                "Relative depth, made metric — and made to fit",
                "The panoramic data")
    eyebrow(s, ML, BODY_TOP, 5.6, "Fitting the response curve")
    bullets(s, ML, BODY_TOP + 0.40, 5.6, [
        "Panoramic depth is **relative**. The photograph's own metric depth "
        "calibrates it.",
        "14 quantile bins, median knots, forced monotonic.",
        "Past the fitted range, extrapolate in log-depth — **capped at 20×**. "
        "Uncapped it invented tens of kilometres.",
    ], size=13.5, gap=15)
    eyebrow(s, 6.98, BODY_TOP, 5.63, "Fitting the terrain budget", AMBER)
    bullets(s, 6.98, BODY_TOP + 0.40, 5.63, [
        "Everything has to fit **D{max} = 100 m**.",
        "A uniform rescale broke Rainier: **6.5×**, and the camera ended up buried "
        "in its own terrain.",
        "So compress only the far end. **The near field stays metric.**",
    ], size=13.5, gap=15, dash_color=AMBER)
    add_pic(s, F("rainier", "depth-calibrated.png"), ML, 4.10, 11.89, 2.08)
    caption(s, ML, 6.30, 11.89,
            "Calibrated panoramic depth, near field bright. Past the treeline the "
            "raw prediction has saturated and the whole range collapses into a "
            "narrow band. Sky is masked out rather than assigned a depth.",
            size=CAP)

    # -- 9 -- height field ----------------------------------------------------
    s = d.slide("Terrain geometry", "Measured ground is never modified", "Terrain")
    panels = [("heightmap-observed.png", "Projected observations",
               "**7%** of the grid, around the camera."),
              ("heightmap-reconstructed.png", "Harmonic completion",
               "Measured cells are boundary conditions."),
              ("heightmap-eroded.png", "Erosion refinement",
               "Only where nothing was observed.")]
    for i, (f, head, body) in enumerate(panels):
        x = ML + i * 2.95
        add_pic(s, F("rainier", f), x, BODY_TOP, 2.55, 2.55)
        tf = textbox(s, x, BODY_TOP + 2.70, 2.72, 1.0)
        style(tf.paragraphs[0].add_run(), 12.5, INK, bold=True).text = head
        p = tf.add_paragraph()
        p.space_before = Pt(5)
        p.line_spacing = 1.2
        rich(p, body, 10.5, MUTED, bold_color=AMBER)
    bullets(s, 9.72, BODY_TOP, 2.89, [
        "A **4096²** grid over 200 m.",
        "Every cell keeps the panorama pixel it was seen through.",
        "The solve never moves a measured cell.",
    ], size=12.5, gap=13)
    line(s, ML, 5.62, W - MR, 5.62, HAIR)
    eyebrow(s, ML, 5.82, 11.9, "De-spoking", CLAY)
    tf = textbox(s, ML, 6.14, 11.9, 0.8)
    p = tf.paragraphs[0]
    p.line_spacing = 1.28
    rich(p, "Per-column depth bias becomes a radial streak — inside the observed "
         "data, where the solve cannot touch it. A tangential high-pass removes "
         "only the angular component: elevation profile unchanged, cliff mask "
         "bit-identical, mean change **0.19 m**.", NOTE, BODY)

    # -- 10 -- ridge anchoring -------------------------------------------------
    s = d.slide("Terrain — the load-bearing assumption",
                "The sky boundary is assumed to be the crest of ground", "Terrain")
    add_pic(s, F("rainier", "silhouette.png"), ML, BODY_TOP, 3.4, 3.4)
    caption(s, ML, BODY_TOP + 3.52, 3.4,
            "The extracted sky silhouette, top-down over the 200 m grid.",
            size=CAP)
    bullets(s, 4.55, BODY_TOP + 0.10, 8.06, [
        "Beyond depth range there is no measurement. The silhouette is the only "
        "signal left.",
        "Each column's first non-sky pixel becomes an elevation, placed at a fixed "
        "radius.",
        "A mountain kilometres away is **brought close and shrunk**: shape is "
        "preserved, distance is not.",
        "Erosion strength follows the same silhouette's jaggedness.",
    ], size=BULLET, gap=20)
    rect(s, 4.55, 5.34, 8.06, 1.06, CLAY, alpha=10, radius=0.06,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 4.55, 5.34, 0.045, 1.06, CLAY)
    tf = textbox(s, 4.92, 5.48, 7.4, 0.85, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.line_spacing = 1.26
    rich(p, "**The assumption: that boundary is the crest of ground.** It holds "
         "for Rainier and Shark Fin. It breaks Paris — slide 18.", 13, BODY,
         bold_color=CLAY)

    # -- 11 -- texturing -------------------------------------------------------
    s = d.slide("Terrain appearance",
                "Projection smears the ground; patterns do not",
                "Terrain appearance")
    eyebrow(s, ML, BODY_TOP, 4.2, "Why projection fails", CLAY)
    bullets(s, ML, BODY_TOP + 0.40, 4.2, [
        "**Grazing incidence.** At 80 m one panorama row covers **12 m** of "
        "ground.",
        "**Upright content.** A 0.5 m flower smears into a streak metres long — a "
        "radial pinwheel from above.",
    ], size=13, gap=16, dash_color=CLAY)
    add_pic(s, F("rainier", "tile-terrain.jpg"), 5.15, BODY_TOP, 3.3, 3.3)
    caption(s, 5.15, BODY_TOP + 3.42, 3.3,
            "The baked terrain layer, top-down in texture space.", size=CAP)
    eyebrow(s, 8.85, BODY_TOP, 3.76, "A layered splat material")
    bullets(s, 8.85, BODY_TOP + 0.40, 3.76, [
        "Real photographic colour only where projection is trustworthy.",
        "Otherwise **real material patches**, cut from the panorama and assembled "
        "over the mesh's triangles.",
        "Generated tiles only where no real reference exists.",
    ], size=13, gap=16)
    line(s, ML, 5.72, W - MR, 5.72, HAIR)
    tf = textbox(s, ML, 5.94, 11.9, 0.6)
    p = tf.paragraphs[0]
    p.line_spacing = 1.28
    rich(p, "Tiles are prompted as an **overhead flat lay under diffuse light**. "
         "The tile is an albedo the client will light — shading baked into it "
         "would be applied twice.", NOTE, MUTED, bold_color=BODY)

    # -- 12 -- objects ---------------------------------------------------------
    s = d.slide("Object decomposition",
                "CLIP proposes; an independent model disposes", "Objects")
    add_pic(s, F("rainier", "clusters.jpg"), ML, BODY_TOP, 7.6, 3.8)
    caption(s, ML, BODY_TOP + 3.92, 7.6,
            "Detections after appearance clustering, coloured by bucket. "
            "Segmentation runs on the panorama, so detections cover the full "
            "360°.", size=CAP)
    bullets(s, 8.68, BODY_TOP, 3.93, [
        "SAM 2 on the panorama, tiled and merged.",
        "CLIP proposes a label; **Grounding DINO has to find it** in the crop.",
        "CLIP always ranks — with no real signal, something still wins.",
        "Sub-clustered by DINOv2 appearance: flower colours, tree species.",
        "Buckets are bounded. An unbounded cut gave **302 groups**.",
    ], size=13, gap=15)
    line(s, ML, 6.16, W - MR, 6.16, HAIR)
    tf = textbox(s, ML, 6.38, 11.9, 0.5)
    rich(tf.paragraphs[0], "One mesh per (class, bucket), instanced at every "
         "member — reconstruction is paid per bucket, not per instance.",
         NOTE, MUTED, bold_color=BODY)

    # -- 13 -- level of detail -------------------------------------------------
    s = d.slide("Level of detail", "Resolvability decides the representation",
                "Objects")
    tiers = [
        ("Mesh", TEAL, "reconstructed once per group, instanced — when it renders "
                       "large enough to resolve as geometry"),
        ("Crossed card", AMBER, "three planes, twelve triangles; holds a "
                                "silhouette from any angle, and inherits the "
                                "sway rig"),
        ("Billboard", SLATE, "one camera-facing quad. Last resort: in stereo it "
                             "reads flat and swings with the head"),
    ]
    for i, (name, col, body) in enumerate(tiers):
        x = ML + i * 4.03
        rect(s, x, BODY_TOP, 3.79, 1.72, PANEL, radius=0.07,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, BODY_TOP, 3.79, 0.035, col)
        tf = textbox(s, x + 0.32, BODY_TOP + 0.18, 3.2, 1.36,
                     anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        style(p.add_run(), 10, MUTED, bold=True).text = f"{i + 1}   "
        style(p.add_run(), 15, col, bold=True).text = name
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.line_spacing = 1.22
        style(p.add_run(), 11.5, BODY).text = body
    bullets(s, ML, 3.80, 6.1, [
        "Crop chosen by **resolution**, not score: too few pixels returns a "
        "different object — a flat sheet.",
        "Representation chosen by **projected angular size**, not distance.",
    ], size=13, gap=15)
    rect(s, 6.98, 3.80, 5.63, 1.66, PANEL, radius=0.06,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = textbox(s, 7.32, 3.98, 5.0, 1.3)
    style(tf.paragraphs[0].add_run(), 10, CLAY, bold=True,
          spacing=140).text = "WHY CARDS EXIST, AND ONLY FOR VEGETATION"
    p = tf.add_paragraph()
    p.space_before = Pt(9)
    p.line_spacing = 1.26
    rich(p, "SAM 3D returns a conifer at **3.50:1** against detections at "
         "**0.12:1** — a 28× disagreement. 70 of 75 instances fall back to cards.",
         12.5, BODY, bold_color=AMBER)
    line(s, ML, 5.72, W - MR, 5.72, HAIR)
    counts = [("128", "meshes"), ("9,303", "cards"), ("33", "billboards"),
              ("2.48 M", "triangles, from 4.47 M")]
    for i, (v, lab) in enumerate(counts):
        x = ML + i * 3.05
        tf = textbox(s, x, 5.94, 2.9, 0.9)
        p = tf.paragraphs[0]
        p.line_spacing = 1.0
        style(p.add_run(), 28, AMBER, bold=True, spacing=-18).text = v
        p = tf.add_paragraph()
        p.space_before = Pt(7)
        style(p.add_run(), 11, MUTED).text = lab

    # -- 14 -- populations -----------------------------------------------------
    s = d.slide("Population synthesis",
                "A detector finds tens of flowers in a meadow of thousands",
                "Population synthesis")
    tf = textbox(s, ML, BODY_TOP, 11.9, 0.5)
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    rich(p, "Place only what was detected and the eye sees **the sampling "
         "pattern**, not the meadow.", 16, BODY)
    add_pic(s, F("rainier", "synthesis-tree.png"), ML, 2.68, 2.5, 2.5)
    add_pic(s, zoom_to_content(F("rainier", "grass-area.png"), pad=0.35),
            7.36, 2.68, 2.5, 2.5)
    blocks = [
        (3.42, 3.0, TEAL, "Population synthesis",
         "*Where else would this class be?* Match the measured pair-correlation "
         "over the rest of the admissible region."),
        (10.10, 2.51, GREEN, "Ground cover",
         "*What is the ground made of?* A mask, not instances — read from the "
         "**original** panorama, which still has its meadow."),
    ]
    for x, w, col, head, body in blocks:
        tf = textbox(s, x, 2.68, w, 2.5)
        style(tf.paragraphs[0].add_run(), 15, col, bold=True).text = head
        p = tf.add_paragraph()
        p.space_before = Pt(8)
        p.line_spacing = 1.28
        rich(p, body, 12.5, BODY, bold_color=INK)
    caption(s, ML, 5.26, 2.5, "Synthesised tree population.", size=CAP,
            align=PP_ALIGN.CENTER)
    caption(s, 7.36, 5.26, 2.5, "Ground-cover mask (zoomed).", size=CAP,
            align=PP_ALIGN.CENTER)
    line(s, ML, 5.66, W - MR, 5.66, HAIR)
    eyebrow(s, ML, 5.88, 11.9, "Colour settles what the label set cannot", AMBER)
    tf = textbox(s, ML, 6.20, 11.9, 0.8)
    p = tf.paragraphs[0]
    p.line_spacing = 1.28
    rich(p, "One ground class folds meadow together with snowfield. The "
         "**excess-green of each cell's own panorama pixel** separates them — then "
         "the same measurement weights the density.", NOTE, BODY)

    # -- 15 -- lighting --------------------------------------------------------
    s = d.slide("Lighting", "Estimate the light, then take it out of everything",
                "Lighting and animation")
    bullets(s, ML, BODY_TOP, 6.6, [
        "LuxDiT recovers an **HDR environment** and a **dominant direction**. The "
        "client uses both.",
        "Every crop already carries the light of the day it was taken. Relighting "
        "it **doubles the illumination**.",
        "IntrinsicDiffusion de-lights crops, so what the client relights is "
        "material colour.",
        "That is also what makes assets **portable** — synthesis puts them "
        "anywhere.",
    ], size=13.5, gap=17)
    add_pic(s, F("rainier", "lighting-ldr.png"), 7.72, BODY_TOP, 2.4, 1.2)
    add_pic(s, F("rainier", "lighting-log.png"), 10.21, BODY_TOP, 2.4, 1.2)
    caption(s, 7.72, BODY_TOP + 1.30, 2.4, "Tone-mapped", size=CAP,
            align=PP_ALIGN.CENTER)
    caption(s, 10.21, BODY_TOP + 1.30, 2.4, "Log scale", size=CAP,
            align=PP_ALIGN.CENTER)
    caption(s, 7.72, BODY_TOP + 1.64, 4.89,
            "The log view exposes the range the client samples: the sun is a "
            "compact lobe well above the sky around it.", size=CAP)
    line(s, ML, 4.62, W - MR, 4.62, HAIR)
    eyebrow(s, ML, 4.82, 11.9, "And the sky itself")
    add_pic(s, F("rainier", "sky.jpg"), ML, 5.16, 11.89, 1.30)
    caption(s, ML, 6.54, 11.89,
            "The sky-inpainted panorama is skybox and image-based lighting source "
            "— one artefact doing two jobs.", size=CAP)

    # -- 16 -- motion ----------------------------------------------------------
    s = d.slide("Animation", "Motion is measured; the wind is authored",
                "Lighting and animation")
    for i, (f, t) in enumerate([("motion-26.jpg", "0 s"), ("motion-54.jpg", "1.2 s"),
                                ("motion-82.jpg", "2.3 s"),
                                ("motion-120.jpg", "3.9 s")]):
        x = ML + i * 1.16
        add_pic(s, F("rainier", f), x, BODY_TOP, 1.02, 1.58)
        caption(s, x, BODY_TOP + 1.66, 1.02, t, size=CAP, align=PP_ALIGN.CENTER)
    caption(s, ML, BODY_TOP + 1.98, 4.5,
            "One tracked flower through the generated video. It tilts and recovers "
            "while staying in place — the signal that types it as anchored.",
            size=CAP)
    bullets(s, 5.34, BODY_TOP, 7.27, [
        "LTX-2 video, **camera locked off** — so drift is the object moving, not "
        "the camera.",
        "Drift under **60% of its own bounding-box diagonal** means stationary.",
        "Stationary gets a **three-bone rig**. Timing is per instance, the "
        "skeleton is shared: a meadow stays one draw.",
        "Anything else is a **rigid body** — velocity and acceleration handed to "
        "physics.",
    ], size=13.5, gap=18)
    line(s, ML, 4.36, 4.9, 4.36, HAIR)
    eyebrow(s, ML, 4.56, 4.6, "One shared wind field")
    tf = textbox(s, ML, 4.90, 4.6, 1.8)
    p = tf.paragraphs[0]
    p.line_spacing = 1.28
    rich(p, "One global field, not per-object oscillators. A gust crosses the "
         "meadow and everything in its path answers in turn.", 12.5, BODY)
    p = tf.add_paragraph()
    p.space_before = Pt(10)
    p.line_spacing = 1.28
    rich(p, "Still authored. Deriving it from the video is future work.", 12.5,
         MUTED)

    # -- 17 -- results ---------------------------------------------------------
    s = d.slide("Results", "Three captures, chosen to break different assumptions",
                "Results")
    tf = textbox(s, ML, BODY_TOP, 11.9, 0.5)
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    rich(p, "No ground-truth 3D exists, so evaluation is **diagnostic**: "
         "instrument, measure, and localise each defect to a stage.", 15, BODY)
    caps = [
        ("rainier-scene.png", "Mount Rainier", TEAL,
         "alpine meadow, dense vegetation, snow",
         "128 / 9,303 / 33", "+35 m", "Conifer reconstruction"),
        ("paris-scene.png", "Paris", AMBER,
         "flat urban river scene, architecture",
         "27 / 4 / 56", "+78 m", "Skyline read as landform"),
        ("sharkfin-scene.png", "Shark Fin Cove", BLUE,
         "coastal cliffs, open water, sea stacks",
         "0 / 1 / 42", "+81 m", "Detections are rock fragments"),
    ]
    for i, (f, name, col, sub, split, relief, limit) in enumerate(caps):
        x = ML + i * 4.03
        add_pic(s, F(f), x, 2.46, 3.79, 2.13)
        tf = textbox(s, x, 4.72, 3.79, 2.1)
        style(tf.paragraphs[0].add_run(), 15, col, bold=True).text = name
        p = tf.add_paragraph()
        p.space_before = Pt(3)
        style(p.add_run(), 10, MUTED).text = sub
        p = tf.add_paragraph()
        p.space_before = Pt(11)
        style(p.add_run(), 11.5, INK, bold=True).text = split
        style(p.add_run(), 10, MUTED).text = "   mesh / card / billboard"
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        style(p.add_run(), 11.5, INK, bold=True).text = relief
        style(p.add_run(), 10, MUTED).text = "   relief added"
        p = tf.add_paragraph()
        p.space_before = Pt(10)
        p.line_spacing = 1.2
        style(p.add_run(), 11, CLAY, bold=True).text = "Breaks on:  "
        style(p.add_run(), 11, BODY).text = limit

    # -- 18 -- failure modes ---------------------------------------------------
    s = d.slide("Results", "Every failure localises to an assumption, not a model",
                "Results")
    a, b = paris_panels()
    if a and b:
        caption(s, ML, BODY_TOP, 4.1, "PARIS HEIGHT FIELD", size=CAP)
        add_pic(s, a, ML, BODY_TOP + 0.32, 1.9, 1.9)
        add_pic(s, b, ML + 2.0, BODY_TOP + 0.32, 1.9, 1.9)
        caption(s, ML, BODY_TOP + 2.34, 4.1,
                "Measured: −6.0 to −0.2 m, correctly flat. After ridge anchoring: "
                "−13.1 to +70.9 m. Hill **A** is the right-bank treeline, **B** "
                "the cathedral. One shared shading scale.", size=CAP, h=1.0)
    fails = [
        ("Skyline read as landform",
         "Terrain-typed silhouette columns: **Paris 0%**, Rainier 98%, Shark Fin "
         "100%. Not one column of Paris's horizon is a landform."),
        ("Confident misclassification on texture",
         "49 detections at **median confidence 0.95**. A rock typed *person*, "
         "captioned “a piece of pizza with bacon on it”."),
        ("Water followed the terrain beneath it",
         "One **11,872 m²** body of sea spanning **73.5 m** of elevation. The "
         "ocean draped up the cliffs."),
        ("Far-field texel density",
         "**67.8%** of Shark Fin's terrain over 0.5 m per texel, against 20.9%. "
         "Within 25 m the captures are indistinguishable."),
    ]
    y = BODY_TOP - 0.02
    for head, body in fails:
        line(s, 5.02, y - 0.16, W - MR, y - 0.16, FAINT)
        tf = textbox(s, 5.02, y, 7.59, 1.2)
        style(tf.paragraphs[0].add_run(), 14, INK, bold=True).text = head
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.line_spacing = 1.26
        rich(p, body, 12, BODY, bold_color=AMBER)
        y += 1.26
    tf = textbox(s, ML, 5.70, 4.1, 0.8)
    p = tf.paragraphs[0]
    p.line_spacing = 1.28
    rich(p, "A defect can surface far from its cause: size compression starts in "
         "depth calibration and shows up **eight stages later**.", 11.5, MUTED,
         bold_color=BODY)

    # -- 19 -- remediation -----------------------------------------------------
    s = d.slide("Remediation",
                "Make the premise testable, not the threshold tunable", "Results")
    tf = textbox(s, ML, BODY_TOP, 11.9, 0.5)
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    rich(p, "Each fix makes an implicit premise **testable from evidence already "
         "computed** — and is measured on every capture, not just the one that "
         "motivated it.", 14, BODY)
    rows = [
        ("Silhouette composition gate", "The sky boundary is the crest of ground",
         "0% / 98% / 100%. Suppressed on Paris only; others bit-identical."),
        ("Per-crop region plausibility",
         "This class can occupy the region its box sits in",
         "Removes 9 / 4 / 12 detections. No false positives."),
        ("Caption–scene plausibility",
         "The caption describes something this scene can contain",
         "A further 2 / 9 / 16. Needs four conjunctive conditions."),
        ("Water surface levelling", "The ground beneath water is flat",
         "Shark Fin 73.5 m spread down to 0.00 m; shoreline step 0.14 m."),
        ("Projected-size level of detail", "Distance predicts resolvability",
         "Rainier 545 meshes down to 296; 4.47 M triangles to 2.48 M."),
        ("Resolution-based crop selection",
         "The best-scoring crop is the best reconstruction input",
         "Crop area up 1.0–9.2×. 205 of 471 instances had been degenerate."),
        ("Geometric UV-fold repair",
         "Adjacent grid cells observed adjacent panorama rows",
         "5,624 triangles collapsed on Rainier (7.0%). Face count unchanged."),
        ("Occluder-footprint inpainting",
         "What foreground removal left behind is ground",
         "Fabricated fill below the horizon: 23.9% down to 17.1%."),
    ]
    cols = [(ML, 3.1), (3.96, 3.7), (7.84, 4.77)]
    y = 2.66
    for lab, (x, w) in zip(("MEASURE", "PREMISE MADE TESTABLE", "MEASURED EFFECT"),
                           cols):
        tf = textbox(s, x, y, w, 0.24)
        style(tf.paragraphs[0].add_run(), 8.5, MUTED, bold=True,
              spacing=120).text = lab
    y += 0.30
    line(s, ML, y, W - MR, y, HAIR)
    y += 0.14
    for i, row in enumerate(rows):
        for text, (x, w), col, bold, ital in zip(
                row, cols, (INK, MUTED, BODY), (True, False, False),
                (False, True, False)):
            tf = textbox(s, x, y, w, 0.5)
            p = tf.paragraphs[0]
            p.line_spacing = 1.18
            style(p.add_run(), 11, col, bold=bold, italic=ital).text = text
        y += 0.49
        if i < len(rows) - 1:
            line(s, ML, y - 0.08, W - MR, y - 0.08, FAINT, 0.5)

    # -- 20 -- conclusions -----------------------------------------------------
    s = d.slide("Conclusions", "Learned models answer *what*; classical methods "
                "decide *how*", "Conclusions and future work")
    bullets(s, ML, BODY_TOP + 0.06, 6.7, [
        "Learned and classical, roughly **two to one** by stage count.",
        "Generative models supply what was never observed. Nothing classical can.",
        "They do not supply a collidable surface, a samplable density, or a mesh "
        "with a skeleton in it.",
        "**12 of 38 active stages use no learned model** — and those are the ones "
        "producing behaviour.",
        "The surviving defects are **assumption failures, not model failures**.",
    ], size=BULLET, gap=19)
    rect(s, 7.72, BODY_TOP - 0.06, 4.89, 4.28, PANEL, radius=0.05,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = textbox(s, 8.06, BODY_TOP + 0.14, 4.25, 3.88, anchor=MSO_ANCHOR.MIDDLE)
    style(tf.paragraphs[0].add_run(), 10, AMBER, bold=True,
          spacing=140).text = "FUTURE WORK"
    for head, body in [
        ("Explicit scene classification", "so no stage assumes a scene type"),
        ("Confidence propagation", "so the system can say where it is guessing"),
        ("Thin-structure reconstruction", "species-aware priors for vegetation"),
        ("Joint refinement", "depth, geometry and placement together"),
        ("Perceptual evaluation", "a user study — our clearest gap"),
    ]:
        p = tf.add_paragraph()
        p.space_before = Pt(13)
        p.line_spacing = 1.22
        hanging(p, 0.14)
        style(p.add_run(), 12.5, INK, bold=True).text = head + "  "
        style(p.add_run(), 12.5, MUTED).text = body
    line(s, ML, 6.14, W - MR, 6.14, HAIR)
    tf = textbox(s, ML, 6.36, 11.9, 0.6)
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    rich(p, "The boundaries are legible and measurable. Addressing them means "
         "**making implicit commitments explicit**.", 15, BODY, bold_color=TEAL)

    # -- 21 -- references ------------------------------------------------------
    s = d.slide("References", "Work this builds on", "References")
    refs = [
        [("Panorama and scene generation", [
            "Kalischek et al. 2025 — CubeDiff",
            "Huang et al. 2025 — DreamCube",
            "Paliwal et al. 2025 — PanoDreamer",
            "Yang et al. 2025 — LayerPano3D",
            "Xie 2025 — WorldGen",
            "Schnepf et al. 2026 — SphericalDreamer",
            "Jin et al. 2026 — SonoWorld",
            "Wu et al. 2026 — 360Anything"]),
         ("Depth estimation", [
            "Lin et al. 2025 — Depth Anything 3",
            "Lin et al. 2025 — Depth Any Panoramas"]),
         ("Image synthesis", [
            "Esser et al. 2024 — Rectified flow transformers",
            "Rombach et al. 2022 — Latent diffusion",
            "Conde et al. 2022 — Swin2SR"])],
        [("Segmentation, detection, recognition", [
            "Ravi et al. 2025 — SAM 2",
            "Liu et al. 2024 — Grounding DINO",
            "Radford et al. 2021 — CLIP",
            "Li et al. 2022 — BLIP",
            "Xiao et al. 2024 — Florence-2",
            "Zhang et al. 2024 — Recognize Anything",
            "Huang et al. 2024 — Tag2Text",
            "Huang et al. 2025 — RAM++",
            "Xie et al. 2021 — SegFormer",
            "Zheng et al. 2024 — BiRefNet",
            "Oquab et al. 2024 — DINOv2"]),
         ("Lighting and video", [
            "Liang et al. 2025 — LuxDiT",
            "Luo et al. 2024 — IntrinsicDiffusion",
            "Lightricks 2026 — LTX-2"])],
        [("Terrain, texture and arrangement", [
            "Jain et al. 2026 — Pixels2Peaks",
            "Génevaux et al. 2015 — Terrain from feature primitives",
            "Barazandeh et al. 2025 — Earthbender",
            "Emilien et al. 2015 — WorldBrush",
            "Barnhart 2020; Hobley 2017; Hutton 2020 — Landlab",
            "Neyret & Cani 1999 — Pattern-based texturing",
            "Geyer & Møller 1994 — Spatial point processes",
            "Hurtut et al. 2009 — Element arrangement by example"]),
         ("Object reconstruction and inpainting", [
            "SAM 3D Team 2025 — SAM 3D",
            "Huang et al. 2025 — SPAR3D",
            "Xiang et al. 2025 — TRELLIS",
            "Suvorov et al. 2022 — LaMa",
            "Li et al. 2022 — MAT",
            "Yu et al. 2023 — Inpaint Anything",
            "Zhao et al. 2026 — ObjectClear",
            "Lee et al. 2024 — Tree-D Fusion"])],
    ]
    accents = [TEAL, AMBER, GREEN]
    for ci, column in enumerate(refs):
        x = ML + ci * 4.03
        y = BODY_TOP
        for gi, (head, items) in enumerate(column):
            tf = textbox(s, x, y, 3.86, 0.24)
            style(tf.paragraphs[0].add_run(), 8.5, accents[ci], bold=True,
                  spacing=110).text = head.upper()
            y += 0.28
            tf = textbox(s, x, y, 3.86, 0.2 * len(items) + 0.3)
            for ii, it in enumerate(items):
                p = para(tf, first=(ii == 0))
                p.line_spacing = 1.16
                p.space_after = Pt(3)
                hanging(p, 0.12)
                style(p.add_run(), 9.5, BODY).text = it
            y += 0.208 * len(items) + 0.30
    line(s, ML, 6.44, W - MR, 6.44, HAIR)
    caption(s, ML, 6.62, 11.9,
            "Full bibliography — 46 entries — in the report.", size=CAP)

    d.prs.save(OUT)
    return d.n


# =============================================================================
#  slide 6: the pipeline, drawn natively so it matches the deck
# =============================================================================
def pipeline_diagram(s):
    y_top, row_h, pitch = 1.78, 0.60, 0.73
    x_pano, w_pano = 2.32, 1.72
    x_track, w_track = 4.44, 6.55
    x_asm, w_asm = 11.34, 1.27
    mid = y_top + (5 * pitch + row_h) / 2

    box = rect(s, ML, mid - 0.41, 1.30, 0.82, PANEL2, radius=0.09,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=HAIR)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.15
    style(p.add_run(), 11, INK, bold=True).text = "Single\nphotograph"
    line(s, ML + 1.30, mid, x_pano, mid, MUTED, 1.0)

    caption(s, x_pano, mid - 1.24, w_pano, "THREE PANORAMA VARIANTS", size=8,
            align=PP_ALIGN.CENTER)
    for i, (name, col) in enumerate([("Original", TEAL),
                                     ("Object-removed", AMBER),
                                     ("Sky-inpainted", BLUE)]):
        chip(s, x_pano, mid - 0.98 + i * 0.62, w_pano, 0.50, name, col, size=10)

    tracks = [
        ("Terrain", ["Height field", "Constrained reconstruction", "Mesh + texture"]),
        ("Objects", ["Detection + typing", "Appearance clustering", "3D assets"]),
        ("Ground cover", ["Spatial statistics", "Population synthesis"]),
        ("Water", ["Region extraction", "Levelled surface"]),
        ("Sky", ["Lighting estimation", "Skybox"]),
        ("Motion", ["Video generation", "Motion extraction", "Rigging"]),
    ]
    badge_w, gap = 1.16, 0.09
    for i, (name, stages) in enumerate(tracks):
        y = y_top + i * pitch
        col = TRACK_COLORS[name]
        rect(s, x_track, y, w_track, row_h, col, alpha=7, radius=0.11,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = textbox(s, x_track + 0.16, y + 0.19, badge_w, 0.3)
        style(tf.paragraphs[0].add_run(), 10, col, bold=True).text = name
        avail = w_track - badge_w - 0.34
        cw = (avail - gap * (len(stages) - 1)) / len(stages)
        for j, st in enumerate(stages):
            cx = x_track + badge_w + 0.22 + j * (cw + gap)
            b = rect(s, cx, y + 0.09, cw, row_h - 0.18, PANEL2, radius=0.16,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=HAIR, line_w=0.5)
            tf = b.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            style(p.add_run(), 9.5, BODY).text = st
            if j:
                line(s, cx - gap, y + row_h / 2, cx, y + row_h / 2, MUTED, 0.75)
        line(s, x_pano + w_pano, mid, x_track, y + row_h / 2, HAIR, 0.75)
        line(s, x_track + w_track, y + row_h / 2, x_asm, mid, HAIR, 0.75)

    box = rect(s, x_asm, mid - 0.41, w_asm, 0.82, PANEL2, radius=0.09,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=TEAL, line_w=1.0)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.15
    style(p.add_run(), 11, INK, bold=True).text = "Scene\nassembly"
    caption(s, x_asm - 0.15, mid + 0.50, w_asm + 0.3,
            "streamed to the client\nas a .frame archive", align=PP_ALIGN.CENTER)

    line(s, ML, 4.90, 4.04, 4.90, HAIR)
    tf = textbox(s, ML, 5.12, 3.5, 1.6)
    p = tf.paragraphs[0]
    p.line_spacing = 1.3
    rich(p, "**Geometry from the object-removed panorama; semantics from the "
         "original.**", 12.5, BODY)
    p = tf.add_paragraph()
    p.space_before = Pt(11)
    p.line_spacing = 1.3
    rich(p, "The server streams assets as they finish. The Unity client carries "
         "the shading, the wind and the physics.", 12.5, MUTED)


if __name__ == "__main__":
    print(f"wrote {OUT} ({build()} slides)")
